import os
from pptx import Presentation
from pptx.util import Inches, Pt
from azure.storage.blob import BlobClient
from typing import List, Dict, Any, Tuple, Optional
import pprint
import tempfile
from urllib.parse import urlparse
import os

# Constants
BLOB_TIMEOUT = 30

from .client import (
    fetch_file,
    download_pptx_to_temp,
    upload_pptx_to_blob,
    cleanup_temp_file
)


def get_powerpoint_layouts(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Extracts and returns information about all available slide layouts in a PowerPoint file.
              # Don't delete output_local_path if that's our result and no dest_url was provided
            if output_local_path and os.path.exists(output_local_path) and dest_url:
                os.remove(output_local_path)
                print(f"Cleaned up temporary output file: {output_local_path}")
                    
async def update_powerpoint_from_blob( Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        List[Dict[str, Any]]: A list of dictionaries containing layout information
    """
    try:
        # Open the PowerPoint file
        presentation = Presentation(pptx_path)
        
        # List to store layout information
        layouts_info = []
        
        # Iterate through all slide masters
        for slide_master in presentation.slide_masters:
            # Iterate through all slide layouts in the slide master
            for idx, slide_layout in enumerate(slide_master.slide_layouts):
                # Get layout name
                layout_name = slide_layout.name
                
                # Get placeholder information
                placeholders = []
                for placeholder in slide_layout.placeholders:
                    placeholders.append({
                        'idx': placeholder.placeholder_format.idx,
                        'type': str(placeholder.placeholder_format.type),
                        'name': placeholder.name if hasattr(placeholder, 'name') else "Unnamed"
                    })
                
                # Add layout information to the list
                layouts_info.append({
                    'index': idx,
                    'name': layout_name,
                    'placeholder_count': len(placeholders),
                    'placeholders': placeholders
                })
                
        return layouts_info
    
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        return []
    
def check_content_fit(content: str, max_title_words: int = 10, max_lines: int = 15, 
                  max_bullets_per_column: int = 8, max_chars_per_line: int = 80) -> dict:
    """
    Checks if the provided content is likely to fit well on a PowerPoint slide.
    Returns warnings if the content might be too large or dense, and suggests font size adjustments.
    
    Args:
        content (str): The content to check
        max_title_words (int): Maximum recommended words for a slide title
        max_lines (int): Maximum recommended total lines for a slide
        max_bullets_per_column (int): Maximum recommended bullet points per column
        max_chars_per_line (int): Maximum recommended characters per line
        
    Returns:
        dict: Dictionary with 'fits' boolean, 'warnings' list, and 'suggested_font_size' dictionary
    """
    lines = content.strip().split('\n')
    total_lines = len(lines)
    warnings = []
    result = {
        'fits': True, 
        'warnings': [],
        'suggested_font_size': {
            'title': None,     # Default font size for title (None means use default)
            'subtitle': None,  # Default font size for subtitle
            'body': None       # Default font size for body text
        },
        'content_complexity': 0  # Measure of overall complexity (0-100)
    }
    
    # Check for empty content
    if not content.strip():
        result['warnings'].append("Content is empty")
        return result
    
    # Initialize complexity score
    complexity_score = 0
    
    # Check title length
    if lines and len(lines[0].split()) > max_title_words:
        title_words = len(lines[0].split())
        warnings.append(f"Title contains {title_words} words, which exceeds the recommended maximum of {max_title_words}")
        complexity_score += min(25, (title_words - max_title_words) * 5)  # Add to complexity score
        
        # Suggest smaller font size for long title
        if title_words > max_title_words + 5:
            result['suggested_font_size']['title'] = 28  # Very long title
        elif title_words > max_title_words:
            result['suggested_font_size']['title'] = 32  # Moderately long title
    
    # Check for very long lines
    long_lines = [line for line in lines if len(line) > max_chars_per_line]
    if long_lines:
        warnings.append(f"Found {len(long_lines)} lines that exceed {max_chars_per_line} characters")
        complexity_score += min(25, len(long_lines) * 5)  # Add to complexity score
        
        # If there are many long lines, suggest smaller body text
        if len(long_lines) > 2:
            result['suggested_font_size']['body'] = 16  # Smaller font for body text with long lines
    
    # Count bullet points
    left_bullets = 0
    right_bullets = 0
    in_left_column = False
    in_right_column = False
    
    for line in lines:
        line_trimmed = line.strip()
        
        # Track column markers
        if "LEFT COLUMN:" in line.upper():
            in_left_column = True
            in_right_column = False
            continue
        elif "RIGHT COLUMN:" in line.upper():
            in_left_column = False
            in_right_column = True
            continue
            
        # Count bullets in each column
        is_bullet = line_trimmed.startswith('-') or line_trimmed.startswith('•')
        is_numbered = len(line_trimmed) > 2 and line_trimmed[0].isdigit() and line_trimmed[1] == '.'
        
        if is_bullet or is_numbered:
            if in_right_column:
                right_bullets += 1
            else:
                left_bullets += 1
    
    # Check total number of bullets
    total_bullets = left_bullets + right_bullets
    if total_bullets > max_bullets_per_column * 2:
        warnings.append(f"Total of {total_bullets} bullet points may be too many for one slide")
        complexity_score += min(25, (total_bullets - max_bullets_per_column * 2) * 2)
        
        # Suggest smaller font size for body text if there are many bullets
        if total_bullets > max_bullets_per_column * 2 + 4:
            result['suggested_font_size']['body'] = 14  # Even smaller for many bullet points
        elif total_bullets > max_bullets_per_column * 2:
            result['suggested_font_size']['body'] = 16  # Smaller for moderately many bullet points
    
    # Check per-column bullet count for two-column layouts
    if in_left_column and in_right_column:  # Two columns were detected
        if left_bullets > max_bullets_per_column:
            warnings.append(f"Left column has {left_bullets} bullet points (recommended: {max_bullets_per_column})")
            complexity_score += min(15, (left_bullets - max_bullets_per_column) * 2)
            
        if right_bullets > max_bullets_per_column:
            warnings.append(f"Right column has {right_bullets} bullet points (recommended: {max_bullets_per_column})")
            complexity_score += min(15, (right_bullets - max_bullets_per_column) * 2)
            
        # Suggest font size based on per-column density
        max_column_bullets = max(left_bullets, right_bullets)
        if max_column_bullets > max_bullets_per_column + 4:
            result['suggested_font_size']['body'] = 14  # Quite dense columns
        elif max_column_bullets > max_bullets_per_column:
            result['suggested_font_size']['body'] = 16  # Moderately dense columns
            
    elif total_bullets > max_bullets_per_column:  # Single column with many bullets
        warnings.append(f"Found {total_bullets} bullet points in a single column (recommended: {max_bullets_per_column})")
        complexity_score += min(20, (total_bullets - max_bullets_per_column) * 2)
        
        # Suggest font size based on bullet density
        if total_bullets > max_bullets_per_column + 4:
            result['suggested_font_size']['body'] = 14  # Very dense single column
        elif total_bullets > max_bullets_per_column:
            result['suggested_font_size']['body'] = 16  # Moderately dense single column
    
    # Check total lines
    if total_lines > max_lines:
        warnings.append(f"Content has {total_lines} lines, which may be too many for one slide (recommended: {max_lines})")
        complexity_score += min(25, (total_lines - max_lines) * 2)
        
        # Suggest smaller subtitle font if needed
        if total_lines > max_lines + 5:
            result['suggested_font_size']['subtitle'] = 20  # Smaller subtitle for very dense slides
    
    # Set overall content complexity
    result['content_complexity'] = min(100, complexity_score)
    
    # Set warnings and fit status
    result['warnings'] = warnings
    if warnings:
        result['fits'] = False
    
    return result

async def get_powerpoint_layouts_from_blob_url(blob_url: str) -> List[Dict[str, Any]]:
    """
    Downloads a PowerPoint file from a blob URL and returns information about all available slide layouts.
    
    Args:
        blob_url (str): Blob URL with SAS token to the PowerPoint file
        
    Returns:
        List[Dict[str, Any]]: A list of dictionaries containing layout information
    """
    local_path = None
    
    try:
        # Download the PowerPoint file
        local_path, error = await download_pptx_to_temp(blob_url)
        if error:
            print(f"Error downloading file: {error}")
            return []
            
        # Use the existing function to get layouts from the local file
        layouts_info = get_powerpoint_layouts(local_path)
        
        return layouts_info
        
    except Exception as e:
        print(f"Error getting PowerPoint layouts from blob URL: {e}")
        import traceback
        traceback.print_exc()
        return []
        
    finally:
        # Clean up temporary file
        if local_path and os.path.exists(local_path):
            try:
                os.remove(local_path)
                print(f"Cleaned up temporary file: {local_path}")
            except Exception as cleanup_error:
                print(f"Warning: Error during cleanup of temporary file: {cleanup_error}")


async def get_layout_names_from_blob_url(blob_url: str) -> List[str]:
    """
    Returns a simple list of layout names available in a PowerPoint file stored at a blob URL.
    
    Args:
        blob_url (str): Blob URL with SAS token to the PowerPoint file
        
    Returns:
        List[str]: A list of layout names
    """
    layouts_info = await get_powerpoint_layouts_from_blob_url(blob_url)
    return [layout['name'] for layout in layouts_info]

async def add_slide_from_blob_url(source_url: str, 
                     content: str, 
                     slide_type: Optional[str] = None, 
                     dest_url: Optional[str] = None, 
                     comments: Optional[str] = None,
                     cleanup_temp_files: bool = True,
                     check_fit: bool = True) -> str:
    """
    Adds a slide to a PowerPoint presentation located at a blob URL with SAS token.
    
    Args:
        source_url (str): Source blob URL with SAS token to the PowerPoint file
        content (str): Content to be added to the slide
        slide_type (str, optional): Exact slide layout name to use. If provided, this layout will be used 
                                  without any automated selection. If None or not found, will auto-select 
                                  the best layout based on content.
        dest_url (str, optional): Destination blob URL with SAS token. If None, returns the file locally
        comments (str, optional): Presenter comments to add to the slide
        cleanup_temp_files (bool, optional): Whether to delete temporary files after processing. Default is True.
        check_fit (bool, optional): Whether to check if content might be too large for the slide and adjust font sizes. Default is True.
        
    Returns:
        str: Path to the saved presentation (local file path or blob URL)
    """
    local_path = None
    output_local_path = None
    
    # Initialize font size adjustment info
    font_size_info = {}
    
    # Check if content might be too large for a slide
    if check_fit:
        fit_result = check_content_fit(content)
        if not fit_result['fits']:
            print("⚠️ Content fit warning: The provided content might not fit well on a slide.")
            for warning in fit_result['warnings']:
                print(f"  - {warning}")
                
            # Get font size suggestions
            font_size_info = fit_result['suggested_font_size']
            
            if any(size is not None for size in font_size_info.values()):
                print("\nAdjusting font sizes to improve fit:")
                if font_size_info['title'] is not None:
                    print(f"  - Title font size: {font_size_info['title']} pt")
                if font_size_info['subtitle'] is not None:
                    print(f"  - Subtitle font size: {font_size_info['subtitle']} pt")
                if font_size_info['body'] is not None:
                    print(f"  - Body text font size: {font_size_info['body']} pt")
            else:
                print("Proceeding with slide creation, but consider revising content.")
    
    try:
        # Clean up the URL and remove any newlines or extra whitespace
        clean_url = source_url.strip().replace('\n', '').replace('\r', '')
        
        # Parse the URL to get the blob name for creating temp file with same name
        parsed_url = urlparse(clean_url)
        blob_name = os.path.basename(parsed_url.path)
        
        # Create a temp file to download the PowerPoint
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            local_path = temp_file.name
            
        print(f"Downloading PowerPoint from blob storage: {blob_name}")
        
        # Download the file from the blob URL using async fetch
        # Download the PowerPoint file
        local_path, error = await download_pptx_to_temp(clean_url)
        if error:
            raise Exception(f"Failed to download PowerPoint file: {error}")
            
        print(f"PowerPoint downloaded to: {local_path}")
        
        # Open the PowerPoint presentation
        presentation = Presentation(local_path)
        
        # Get layout information
        layouts_info = get_powerpoint_layouts(local_path)
        
        # Determine which layout to use - prioritize the provided slide_type
        if slide_type:
            # Use the specified slide type directly
            selected_layout_name = slide_type
            print(f"Using specified layout: {selected_layout_name}")
            
            # Check if the specified layout exists (just for warning)
            if not any(layout['name'] == slide_type for layout in layouts_info):
                print(f"Warning: Specified layout '{slide_type}' not found in presentation. Will attempt to use it anyway.")
        else:
            # Auto-select layout based on content only if no layout was specified
            #selected_layout_name = suggest_layout_with_structured_output(content, layouts_info)
            print(f"Auto-selected layout: {selected_layout_name}")
        
        # Find the layout object by name
        layout_obj = None
        
        # First, try exact match
        for master in presentation.slide_masters:
            for layout in master.slide_layouts:
                if layout.name == selected_layout_name:
                    layout_obj = layout
                    break
            if layout_obj:
                break
        
        # If exact match not found and slide_type was manually specified, try case-insensitive match
        if not layout_obj and slide_type:
            selected_layout_lower = selected_layout_name.lower()
            for master in presentation.slide_masters:
                for layout in master.slide_layouts:
                    if layout.name.lower() == selected_layout_lower:
                        layout_obj = layout
                        print(f"Found case-insensitive match for '{selected_layout_name}': '{layout.name}'")
                        break
                if layout_obj:
                    break
                
        # Final fallback to the default layout
        if not layout_obj:
            print(f"Error: Could not find layout object for '{selected_layout_name}'. Using default layout.")
            layout_obj = presentation.slide_layouts[0]
            
        # Add a new slide with the selected layout
        new_slide = presentation.slides.add_slide(layout_obj)
        
        # Add content to the slide based on placeholder types
        if isinstance(content, bytes):
            content = content.decode('utf-8')
        
        lines = content.split('\n')
        title_set = False
        content_added = False
        
        # Find all placeholders and their types
        placeholders = []
        for shape in new_slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                placeholders.append({
                    'shape': shape,
                    'idx': shape.placeholder_format.idx,
                    'type': shape.placeholder_format.type,
                    'name': shape.name if hasattr(shape, 'name') else "Unnamed"
                })
        
        # First, try to find title placeholder (type 1)
        title_placeholders = [p for p in placeholders if p['type'] == 1]
        if title_placeholders and lines:
            title_shape = title_placeholders[0]['shape']
            title_shape.text = lines[0]
            
            # Apply font size adjustment for title if needed
            if check_fit and font_size_info.get('title') is not None:
                from pptx.util import Pt
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size_info['title'])
                        
            lines = lines[1:]
            title_set = True
            
        # Then, try to find subtitle placeholder (type 2) if it exists and we have more content
        subtitle_placeholders = [p for p in placeholders if p['type'] == 2]
        if subtitle_placeholders and lines and len(lines[0].strip()) < 50:  # Only use short lines for subtitles
            subtitle_shape = subtitle_placeholders[0]['shape']
            subtitle_shape.text = lines[0]
            
            # Apply font size adjustment for subtitle if needed
            if check_fit and font_size_info.get('subtitle') is not None:
                from pptx.util import Pt
                for paragraph in subtitle_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size_info['subtitle'])
                        
            lines = lines[1:]
            
        # Function to apply body text font size to a shape
        def apply_body_text_font_size(shape):
            if check_fit and font_size_info.get('body') is not None:
                from pptx.util import Pt
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size_info['body'])
        
        # Handle content placeholders - look for body/content placeholders
        content_placeholder_types = [3, 7]  # Common content placeholder types
        content_placeholders = [p for p in placeholders if p['type'] in content_placeholder_types]
        
        # Special handling for comparison/two-column layouts
        left_column_content = []
        right_column_content = []
        
        # Check if content contains explicit LEFT/RIGHT COLUMN markers
        column_mode = False
        current_column = "left"
        
        for line in lines:
            if "LEFT COLUMN:" in line.upper():
                column_mode = True
                current_column = "left"
                continue
            elif "RIGHT COLUMN:" in line.upper():
                column_mode = True
                current_column = "right"
                continue
                
            if column_mode:
                if current_column == "left":
                    left_column_content.append(line)
                else:
                    right_column_content.append(line)
            else:
                # If no explicit columns, add to left column for now
                left_column_content.append(line)
                
        # If we have a two-column layout and content has been split
        if len(content_placeholders) == 2 and column_mode:
            content_placeholders[0]['shape'].text = '\n'.join(left_column_content)
            apply_body_text_font_size(content_placeholders[0]['shape'])
            
            content_placeholders[1]['shape'].text = '\n'.join(right_column_content)
            apply_body_text_font_size(content_placeholders[1]['shape'])
            
            content_added = True
        # If two column layout but no explicit column markers, split content approximately in half
        elif len(content_placeholders) == 2 and not column_mode:
            middle = len(left_column_content) // 2
            right_column_content = left_column_content[middle:]
            left_column_content = left_column_content[:middle]
            
            content_placeholders[0]['shape'].text = '\n'.join(left_column_content)
            apply_body_text_font_size(content_placeholders[0]['shape'])
            
            content_placeholders[1]['shape'].text = '\n'.join(right_column_content)
            apply_body_text_font_size(content_placeholders[1]['shape'])
            
            content_added = True
        # If not a two-column layout or we couldn't distribute content
        elif content_placeholders and not content_added:
            # Put all remaining content in the first content placeholder
            content_placeholders[0]['shape'].text = '\n'.join(left_column_content)
            apply_body_text_font_size(content_placeholders[0]['shape'])
            
            content_added = True
            
        # If no suitable content placeholder found but we have content, add a textbox
        if not content_added and (left_column_content or right_column_content):
            from pptx.util import Inches
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(4)
            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            
            # Combine all remaining content
            all_content = left_column_content + right_column_content
            textbox.text_frame.text = '\n'.join(all_content)
            
            # Apply font size adjustment for textbox
            apply_body_text_font_size(textbox)
            
        # Add presenter comments if provided
        if comments:
            # python-pptx doesn't directly support slide comments/notes
            # Instead, we add the comments to the slide notes
            if not hasattr(new_slide, 'notes_slide'):
                print("Warning: This version of python-pptx doesn't support slide notes")
            else:
                notes_slide = new_slide.notes_slide
                notes_slide.notes_text_frame.text = comments
        
        # Create output temp file for saving
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as output_temp:
            output_local_path = output_temp.name
        
        # Save the modified presentation
        presentation.save(output_local_path)
        print(f"Presentation saved to temporary file: {output_local_path}")
        
        # Handle destination URL if provided
        result_path = output_local_path
        if dest_url:
            try:
                # Parse the destination URL
                parsed_dest = urlparse(dest_url)
                if not parsed_dest.netloc:
                    raise ValueError("Invalid destination URL")
                
                # Set up blob client
                blob_client = BlobClient.from_blob_url(dest_url, timeout=BLOB_TIMEOUT)
                
                # Upload the file to blob storage
                with open(output_local_path, "rb") as data:
                    blob_client.upload_blob(data, overwrite=True)
                    
                print(f"Presentation uploaded to: {dest_url}")
                result_path = dest_url
            except Exception as upload_error:
                print(f"Error uploading to blob storage: {upload_error}")
                print(f"File is available locally at: {output_local_path}")
                result_path = output_local_path
        else:
            # Return the local path if no destination URL is provided
            result_path = output_local_path
            
        return result_path
            
    except Exception as e:
        print(f"Error adding slide from blob URL: {e}")
        import traceback
        traceback.print_exc()
        return None
        
    finally:
        # Clean up temporary files if requested
        if cleanup_temp_files:
            if local_path and os.path.exists(local_path):
                try:
                    os.remove(local_path)
                    print(f"Cleaned up temporary source file: {local_path}")
                except Exception as cleanup_error:
                    print(f"Warning: Error cleaning up source file: {cleanup_error}")
                
            # Don't delete output_local_path if that's our result and no dest_url was provided
            if output_local_path and os.path.exists(output_local_path) and dest_url:
                try:
                    os.remove(output_local_path)
                    print(f"Cleaned up temporary output file: {output_local_path}")
                except Exception as cleanup_error:
                    print(f"Warning: Error cleaning up output file: {cleanup_error}")
                    
async def update_powerpoint_from_blob(
    source_url: str, 
    dest_url: str,
    slide_updates: List[Dict[str, Any]]
) -> Tuple[bool, str]:
    """
    Downloads a PowerPoint from a source blob URL, updates it according to provided instructions,
    and uploads it to a destination blob URL.
    
    Args:
        source_url (str): Source blob URL with SAS token where the template PowerPoint is stored
        dest_url (str): Destination blob URL with SAS token where the updated PowerPoint should be saved
        slide_updates (List[Dict[str, Any]]): List of slide update instructions
        
    Returns:
        Tuple[bool, str]: Success flag and error message (if any)
    """
    local_path = None
    success = False
    error_msg = ""
    
    try:
        # Download the source PowerPoint file
        local_path, error = await download_pptx_to_temp(source_url)
        if error:
            return False, f"Error downloading file: {error}"
            
        print(f"PowerPoint downloaded to: {local_path}")
        
        # Update the PowerPoint file
        success = update_powerpoint_file(local_path, slide_updates)
        if not success:
            return False, "Failed to update PowerPoint file"
        
        print(f"PowerPoint updated successfully")
        
        # Upload the updated file to the destination blob
        print(f"Uploading updated PowerPoint to blob storage")
        error = await upload_pptx_to_blob(local_path, dest_url)
        if error:
            return False, f"Error uploading file: {error}"
            
        print(f"PowerPoint uploaded successfully to destination blob")
        success = True
        return True, ""
        
    except Exception as e:
        error_msg = f"Error updating PowerPoint from blob: {str(e)}"
        print(error_msg)
        import traceback
        traceback.print_exc()
        return False, error_msg
        
    finally:
        # Clean up temporary file
        await cleanup_temp_file(local_path)

def update_powerpoint_file(file_path: str, slide_updates: List[Dict[str, Any]]) -> bool:
    """
    Updates a PowerPoint file with the provided slide updates.
    
    Args:
        file_path (str): Path to the PowerPoint file
        slide_updates (List[Dict[str, Any]]): List of slide update instructions
        
    Returns:
        bool: True if update was successful, False otherwise
    """
    try:
        # Open the PowerPoint file
        prs = Presentation(file_path)
        
        # Process each slide update
        for update in slide_updates:
            slide_index = update.get('slide_index')
            if slide_index is None:
                print("Warning: Slide index not provided in update instruction")
                continue
                
            # Make sure slide index is valid
            if slide_index < 0 or slide_index >= len(prs.slides):
                print(f"Warning: Invalid slide index {slide_index}")
                continue
            
            slide = prs.slides[slide_index]
            
            # Process placeholder updates
            updates = update.get('updates', [])
            for placeholder_update in updates:
                idx = placeholder_update.get('placeholder_idx')
                new_text = placeholder_update.get('text')
                font_size = placeholder_update.get('font_size')
                
                if idx is None or new_text is None:
                    print("Warning: Missing placeholder_idx or text in update")
                    continue
                    
                # Find the shape with matching placeholder index
                shape = None
                for s in slide.shapes:
                    if (hasattr(s, 'placeholder_format') and 
                        s.placeholder_format.idx == idx):
                        shape = s
                        break
                        
                if shape is None:
                    print(f"Warning: No placeholder found with index {idx}")
                    continue
                    
                # Update text
                shape.text = new_text
                
                # Update font size if specified
                if font_size:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
                            
        # Save the updated file
        prs.save(file_path)
        return True
        
    except Exception as e:
        print(f"Error updating PowerPoint file: {e}")
        import traceback
        traceback.print_exc()
        return False